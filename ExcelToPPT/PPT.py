#!/usr/bin/env python
# -*- coding: utf-8 -*-
# Author : Dong-Qing 
# Time : 2018/3/7

import os
from win32com.client import Dispatch
from pptx import Presentation

class Pptoperating(object):


    def __init__(self):
        self.pptDir = ""
        self.imgDir = ""
        self.excelDir = ""
        self.name = ""


    def export(self):
        # ppt= Dispatch("PowerPoint.application")
        # ppt.Visible = False
        # wb = ppt.Workbooks.Open(os.path.join(self.pptDir ,self.name))
        # # print(wb)
        ppt_dir = os.path.join(self.pptDir ,self.name)

        # with open(ppt_dir) as f:
        #     source_stream = StringIO(f.read())
        # prs = Presentation(source_stream)
        # source_stream.close()

        prs = Presentation(self.name)

        title_slide_layout = prs.slide_layouts[0]
        f = str(title_slide_layout)

        bb=1

        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Hello,World!"
        subtitle.text = "python-pptx was here!"
        prs.save('test.pptx')

        x= 1

if __name__=='__main__':
    pl = Pptoperating()
    pl.pptDir = r"G:\youyu\ExcelToPPT\ppt" + '\\'
    pl.name = '金鸡项目建设进度报告-尚晓晗2018.3.2.ppt'
    pl.excelDir = r'G:\youyu\ExcelToPPT\excel'
    pl.export()
